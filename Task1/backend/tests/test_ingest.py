"""Ingestion tests that need no network: PPTX round-trip, HTML parsing, URL parsing."""

from __future__ import annotations

import asyncio
from io import BytesIO

import pytest
from pptx import Presentation

from app.ingest.base import IngestError
from app.ingest.pptx import ingest_pptx
from app.ingest.web import _extract, normalise_url
from app.ingest.youtube import extract_video_id, format_timestamp, is_youtube_url


def build_deck() -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    first = presentation.slides.add_slide(layout)
    first.shapes.title.text = "Retrieval Augmented Generation"
    first.placeholders[1].text = "Chunk the document\nEmbed each chunk"
    first.notes_slide.notes_text_frame.text = "Mention that chunk size is a trade-off."

    second = presentation.slides.add_slide(layout)
    second.shapes.title.text = "Evaluation"
    second.placeholders[1].text = "Measure grounding and citation accuracy"

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_pptx_extracts_titles_bodies_and_notes() -> None:
    result = asyncio.run(ingest_pptx(build_deck(), "deck.pptx"))

    assert result.title == "Retrieval Augmented Generation"
    assert [segment.locator for segment in result.segments] == ["slide 1", "slide 2"]
    assert "Chunk the document" in result.segments[0].text
    assert "Speaker notes:" in result.segments[0].text
    assert "citation accuracy" in result.segments[1].text


def test_pptx_rejects_non_pptx_bytes() -> None:
    with pytest.raises(IngestError, match="PowerPoint"):
        asyncio.run(ingest_pptx(b"not a presentation", "broken.pptx"))


HTML = """
<html><head><title>Vector Search Guide</title></head>
<body>
  <nav>Home Docs Blog</nav>
  <main>
    <h1>Vector Search Guide</h1>
    <p>Vector search finds documents by meaning rather than exact keywords, which helps with paraphrase.</p>
    <h2>Installation</h2>
    <p>Install the client library and configure the API key before indexing your first collection.</p>
    <h2>Usage</h2>
    <p>Call the query endpoint with an embedding vector to retrieve the nearest neighbours by cosine.</p>
  </main>
  <footer>Copyright</footer>
</body></html>
"""


def test_web_extraction_uses_headings_as_locators() -> None:
    result = _extract(HTML, "https://example.com/guide")
    locators = [segment.locator for segment in result.segments]

    assert result.title == "Vector Search Guide"
    assert "Installation" in locators
    assert "Usage" in locators
    # Navigation and footer chrome is dropped.
    assert all("Copyright" not in segment.text for segment in result.segments)


def test_web_extraction_rejects_pages_without_prose() -> None:
    with pytest.raises(IngestError, match="No readable article text"):
        _extract("<html><body><div>hi</div></body></html>", "https://example.com")


def test_normalise_url_rejects_non_http_schemes() -> None:
    with pytest.raises(IngestError, match="http"):
        normalise_url("file:///etc/passwd")


def test_normalise_url_adds_https_when_missing() -> None:
    assert normalise_url("example.com/docs").startswith("https://example.com")


@pytest.mark.parametrize(
    "url",
    [
        "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        "https://youtu.be/dQw4w9WgXcQ",
        "https://www.youtube.com/shorts/dQw4w9WgXcQ",
        "https://www.youtube.com/watch?list=PL123&v=dQw4w9WgXcQ",
    ],
)
def test_youtube_url_forms_are_recognised(url: str) -> None:
    assert is_youtube_url(url)
    assert extract_video_id(url) == "dQw4w9WgXcQ"


def test_non_youtube_url_is_treated_as_a_webpage() -> None:
    assert not is_youtube_url("https://example.com/watch?v=abc")


def test_timestamp_formatting() -> None:
    assert format_timestamp(195.4) == "03:15"
    assert format_timestamp(3725) == "1:02:05"
