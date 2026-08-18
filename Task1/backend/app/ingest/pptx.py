"""PPTX ingestion.

Slides carry structure that plain text extraction throws away, so each slide keeps
its title, body text in reading order, table contents, and speaker notes. Notes in
particular often hold the explanation a learner is actually asking about.
"""

from __future__ import annotations

import asyncio
import zipfile
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..schemas import Segment, SourceKind
from .base import IngestError, IngestResult, ensure_not_empty


def _shape_text(shape) -> str:  # noqa: ANN001 - python-pptx shapes are untyped
    parts: list[str] = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs).strip()
            if line:
                parts.append(line)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _slide_title(slide) -> str | None:  # noqa: ANN001
    placeholder = getattr(slide.shapes, "title", None)
    if placeholder is not None and placeholder.has_text_frame:
        title = placeholder.text_frame.text.strip()
        if title:
            return title
    return None


def _extract(data: bytes, filename: str) -> IngestResult:
    try:
        presentation = Presentation(BytesIO(data))
    except (PackageNotFoundError, zipfile.BadZipFile, KeyError, ValueError) as exc:
        # A .pptx is a zip archive; corrupt or misnamed files surface as zip errors.
        raise IngestError(
            "This file could not be read as a PowerPoint (.pptx) file. Note that the older "
            ".ppt format is not supported -- re-save it as .pptx."
        ) from exc

    segments: list[Segment] = []
    deck_title: str | None = None

    for number, slide in enumerate(presentation.slides, start=1):
        title = _slide_title(slide)
        if number == 1:
            deck_title = title

        body = [_shape_text(shape) for shape in slide.shapes]
        lines = [text for text in body if text]

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"Speaker notes: {notes}")

        if not lines and not title:
            continue

        heading = f"{title}\n" if title else ""
        segments.append(
            Segment(
                text=f"{heading}{chr(10).join(lines)}",
                position=number,
                locator=f"slide {number}",
            )
        )

    title = (deck_title or presentation.core_properties.title or Path(filename).stem).strip()
    result = IngestResult(
        kind=SourceKind.PPTX,
        title=title or filename,
        segments=segments,
    )
    return ensure_not_empty(
        result,
        "No text was found in this presentation. Slides made entirely of images "
        "cannot be read without OCR.",
    )


async def ingest_pptx(data: bytes, filename: str) -> IngestResult:
    return await asyncio.to_thread(_extract, data, filename)
